#!/usr/bin/env python3
#
# Optional system dependencies (beyond `pip install -r requirements.txt`):
#   - .rar archives: the 'unrar' (or 'unar'/'bsdtar') tool must be on
#     PATH. Windows: winget install RARLab.WinRAR
#   - Scanned PDFs (OCR): the Tesseract OCR engine must be installed.
#     Windows: winget install UB-Mannheim.TesseractOCR
# Both are optional: without them, the affected archive/PDF is reported
# as an error and the rest of the run continues normally.
#
# Known limitation: PII detection is regex-based and heuristic. In
# particular, four-part version strings ("1.0.0.0" in .cs/.json/.xml or
# build config files) are indistinguishable from IPv4 addresses and are
# redacted as [IP].
#
# A file, archive, or PDF page that fails to process (including PDF
# table detection taking longer than PDF_TABLE_TIMEOUT_SECONDS) is
# reported as an error and EXCLUDED from the output, never copied in
# unprocessed - the output directory must never contain a raw file that
# was supposed to be anonymized but wasn't.

import argparse
import concurrent.futures
import json
import re
import shutil
import sys
import tempfile
import time
from pathlib import Path
from typing import NamedTuple, Optional

# ------------------------------------------------------------
# Supported file types
# ------------------------------------------------------------

TEXT_EXTENSIONS = {
    ".txt", ".md", ".markdown", ".rst",
    ".cs", ".csx",
    ".py", ".pyw",
    ".js", ".jsx", ".ts", ".tsx", ".mjs", ".cjs",
    ".html", ".htm", ".css", ".scss", ".sass", ".less",
    ".java", ".kt", ".kts", ".scala",
    ".c", ".h", ".cpp", ".hpp", ".cc", ".cxx",
    ".go", ".rs", ".rb", ".php", ".swift", ".dart",
    ".fs", ".fsx", ".vb",
    ".json", ".xml", ".yaml", ".yml", ".toml",
    ".ini", ".cfg", ".conf", ".env",
    ".sql", ".sh", ".bash", ".zsh", ".ps1", ".bat", ".cmd",
    ".dockerfile",
}

OFFICE_EXTENSIONS = {".docx", ".xlsx", ".pptx"}


# ------------------------------------------------------------
# Replacement map
# ------------------------------------------------------------

def load_replacements(path: Path) -> dict:
    with path.open("r", encoding="utf-8") as f:
        replacements = json.load(f)

    if not isinstance(replacements, dict):
        raise ValueError("replacements.json must contain a JSON object.")

    # Longer strings first. This prevents a shorter replacement
    # from being applied before a more specific one.
    return dict(
        sorted(
            replacements.items(),
            key=lambda item: len(item[0]),
            reverse=True,
        )
    )


def anonymize_text(text: str, replacements: dict) -> str:
    if not replacements:
        return text

    # Matching is case-insensitive (a term may appear as "BDR" in one
    # file and "bdr" in a Java package name in another), but the
    # substituted value is always the dictionary's exact value, never
    # case-adjusted to match what was found.
    pattern = re.compile(
        "|".join(re.escape(original) for original in replacements),
        re.IGNORECASE,
    )
    lookup = {original.lower(): value for original, value in replacements.items()}

    return pattern.sub(
        lambda match: lookup[match.group(0).lower()],
        text,
    )


# ------------------------------------------------------------
# PII checksum validators
# ------------------------------------------------------------

def validate_pesel(number: str) -> bool:
    if not re.fullmatch(r"\d{11}", number):
        return False

    weights = (1, 3, 7, 9, 1, 3, 7, 9, 1, 3)
    checksum = sum(
        int(digit) * weight
        for digit, weight in zip(number, weights)
    )
    check_digit = (10 - checksum % 10) % 10

    return check_digit == int(number[10])


def validate_nip(number: str) -> bool:
    if not re.fullmatch(r"\d{10}", number):
        return False

    weights = (6, 5, 7, 2, 3, 4, 5, 6, 7)
    checksum = sum(
        int(digit) * weight
        for digit, weight in zip(number, weights)
    )
    check_digit = checksum % 11

    if check_digit == 10:
        return False

    return check_digit == int(number[9])


def validate_iban(code: str) -> bool:
    code = code.replace(" ", "").upper()

    if not re.fullmatch(r"[A-Z]{2}\d{2}[A-Z0-9]{10,30}", code):
        return False

    rearranged = code[4:] + code[:4]

    try:
        digits = "".join(str(int(char, 36)) for char in rearranged)
    except ValueError:
        return False

    return int(digits) % 97 == 1


def validate_luhn(number: str) -> bool:
    if not re.fullmatch(r"\d{13,19}", number):
        return False

    total = 0

    for index, digit in enumerate(reversed(number)):
        value = int(digit)

        if index % 2 == 1:
            value *= 2
            if value > 9:
                value -= 9

        total += value

    return total % 10 == 0


# ------------------------------------------------------------
# PII detection
# ------------------------------------------------------------

EMAIL_PATTERN = re.compile(r"\b[\w.+-]+@[\w-]+\.[\w.-]+\b")

# Only a literal space or dash may separate the digit groups. Using \s
# here would let a "phone number" span newlines and collapse several
# unrelated lines of a document into a single [PHONE] token.
PHONE_PATTERN = re.compile(
    r"(?<!\d)(?:\+48[ -]?)?\d{3}[ -]?\d{3}[ -]?\d{3}(?!\d)"
)

# IBANs are printed either contiguously ("PL6110901014...", sometimes
# lowercase) or - far more commonly - grouped in fours
# ("PL61 1090 1014 0000 0712 1981 2874"). Both forms must be redacted as
# one span; otherwise CARD_PATTERN grabs the 16 digits in the middle and
# the country code and the tail of the account number leak in cleartext.
#
# No fixed-shape regex can find the end of a grouped IBAN: its own
# groups and a short word or number following it are both just "four
# alphanumerics". Polish (28 characters) and Spanish (24) IBANs - this
# tool's main target - have a length that is an exact multiple of four,
# so any trailing token looks exactly like one more group of the number.
#
# Hence this pattern deliberately over-matches: a generous candidate
# span that always contains the whole IBAN (ISO 13616 caps it at 34
# characters) plus some trailing text, and the checksum - not the
# regex - decides where the IBAN really ends. See redact_ibans().
IBAN_CANDIDATE_PATTERN = re.compile(
    r"\b[A-Za-z]{2}\d{2}"
    r"(?:[A-Za-z0-9]{10,30}|(?:[ ][A-Za-z0-9]{1,4}){2,10})"
    r"\b"
)

# Shortest string validate_iban() can accept (its own shape check is
# 2 + 2 + at least 10 characters; the shortest registered IBAN, Norway's,
# is 15). Prefixes below this cannot validate, so the trim loop stops.
IBAN_MIN_LENGTH = 14

CARD_PATTERN = re.compile(
    r"(?<!\d)(?:\d{4}[ -]?){3}\d{4}(?!\d)|(?<!\d)\d{13,19}(?!\d)"
)

PESEL_PATTERN = re.compile(r"(?<!\d)\d{11}(?!\d)")

NIP_PATTERN = re.compile(r"(?<!\d)\d{10}(?!\d)")

# Known, accepted trade-off: this also matches four-part version strings
# such as assembly/package versions ("1.0.0.0", "2.1.0.3") in .cs/.json/
# .xml/build files. Excluding them risks missing real addresses, so the
# false positive is accepted rather than pattern-matched away.
IPV4_PATTERN = re.compile(
    r"\b(?:(?:25[0-5]|2[0-4]\d|1?\d?\d)\.){3}(?:25[0-5]|2[0-4]\d|1?\d?\d)\b"
)

# These categories all share the same shape - match, validate the whole
# match, replace it wholesale - and are applied in order, around the
# IBAN step that runs between them (see redact_pii).
PII_PATTERNS_BEFORE_IBAN = (
    ("EMAIL", EMAIL_PATTERN, None),
    ("IP", IPV4_PATTERN, None),
)

PII_PATTERNS_AFTER_IBAN = (
    (
        "CARD",
        CARD_PATTERN,
        lambda raw: validate_luhn(re.sub(r"[ -]", "", raw)),
    ),
    ("PESEL", PESEL_PATTERN, validate_pesel),
    ("NIP", NIP_PATTERN, validate_nip),
    ("PHONE", PHONE_PATTERN, None),
)


def redact_ibans(text: str) -> tuple[str, int]:
    """
    Replace every checksum-valid IBAN in `text` with "[IBAN]".

    IBAN_CANDIDATE_PATTERN over-matches on purpose, so each candidate is
    trimmed from the right one space-separated group at a time and the
    longest prefix that passes validate_iban() wins. That prefix - and
    only it - is replaced; whatever followed it stays exactly as it was,
    because it is ordinary text, not part of the account number. A
    candidate with no valid prefix was never an IBAN and is left alone.

    Scanning therefore cannot use re.sub(), which would resume after the
    whole over-matched span and skip past text that was never redacted:
    a second IBAN later on the same line, or one standing behind a
    candidate that failed, would be swallowed by the over-match and leak
    in cleartext. Instead the scan resumes right after whatever was
    actually consumed.
    """
    count = 0
    chunks = []
    pos = 0

    while (match := IBAN_CANDIDATE_PATTERN.search(text, pos)) is not None:
        tokens = match.group(0).split(" ")
        prefix = None

        for stop in range(len(tokens), 0, -1):
            head = " ".join(tokens[:stop])

            if len(head.replace(" ", "")) < IBAN_MIN_LENGTH:
                # Every remaining prefix is shorter still.
                break

            if validate_iban(head):
                prefix = head
                break

        chunks.append(text[pos:match.start()])

        if prefix is None:
            # Not an IBAN. Keep it verbatim and resume after its first
            # token: a real IBAN can only start on a token boundary, so
            # nothing is skipped and each step still moves forward.
            pos = match.start() + len(tokens[0])
            chunks.append(text[match.start():pos])
        else:
            chunks.append("[IBAN]")
            count += 1
            pos = match.start() + len(prefix)

    chunks.append(text[pos:])

    return "".join(chunks), count


def _redact_categories(text: str, categories) -> tuple[str, int]:
    total = 0

    for label, pattern, validator in categories:
        def replace(match, label=label, validator=validator):
            nonlocal total

            if validator is not None and not validator(match.group(0)):
                return match.group(0)

            total += 1
            return f"[{label}]"

        text = pattern.sub(replace, text)

    return text, total


def redact_pii(text: str) -> tuple[str, int]:
    text, before = _redact_categories(text, PII_PATTERNS_BEFORE_IBAN)

    # IBANs need their own validation-aware pass (see redact_ibans), and
    # it has to run here - before CARD/PESEL/NIP, which would otherwise
    # grab a still-unredacted IBAN's digit groups, mislabel them and
    # leak the rest of the account number around them.
    text, ibans = redact_ibans(text)

    text, after = _redact_categories(text, PII_PATTERNS_AFTER_IBAN)

    return text, before + ibans + after


# ------------------------------------------------------------
# Office -> Markdown
# ------------------------------------------------------------

def convert_docx_to_markdown(source: Path) -> str:
    try:
        from docx import Document
    except ImportError:
        raise RuntimeError(
            "Missing python-docx. Install it with: pip install python-docx"
        )

    document = Document(source)
    result = []

    for paragraph in document.paragraphs:
        text = paragraph.text.strip()

        if not text:
            result.append("")
            continue

        style = paragraph.style.name.lower()

        if "heading 1" in style:
            result.append(f"# {text}")
        elif "heading 2" in style:
            result.append(f"## {text}")
        elif "heading 3" in style:
            result.append(f"### {text}")
        elif "heading 4" in style:
            result.append(f"#### {text}")
        elif "list bullet" in style:
            result.append(f"- {text}")
        elif "list number" in style:
            result.append(f"1. {text}")
        else:
            result.append(text)

    for table in document.tables:
        result.append("")
        rows = []

        for row in table.rows:
            rows.append([
                cell.text.replace("\n", " ").strip()
                for cell in row.cells
            ])

        if rows:
            result.append("| " + " | ".join(rows[0]) + " |")
            result.append("| " + " | ".join("---" for _ in rows[0]) + " |")

            for row in rows[1:]:
                result.append("| " + " | ".join(row) + " |")

    return "\n".join(result)


def convert_xlsx_to_markdown(source: Path) -> str:
    try:
        from openpyxl import load_workbook
    except ImportError:
        raise RuntimeError(
            "Missing openpyxl. Install it with: pip install openpyxl"
        )

    workbook = load_workbook(source, data_only=True)
    result = []

    for worksheet in workbook.worksheets:
        result.append(f"# {worksheet.title}")
        result.append("")

        rows = list(worksheet.iter_rows(values_only=True))

        while rows and all(value is None for value in rows[-1]):
            rows.pop()

        if not rows:
            continue

        max_columns = max(len(row) for row in rows)

        while max_columns > 0:
            if any(
                len(row) >= max_columns
                and row[max_columns - 1] is not None
                for row in rows
            ):
                break
            max_columns -= 1

        rows = [list(row[:max_columns]) for row in rows]

        def value_to_string(value):
            if value is None:
                return ""
            return (
                str(value)
                .replace("\n", " ")
                .replace("|", r"\|")
            )

        rows = [
            [value_to_string(value) for value in row]
            for row in rows
        ]

        if rows:
            result.append("| " + " | ".join(rows[0]) + " |")
            result.append("| " + " | ".join("---" for _ in rows[0]) + " |")

            for row in rows[1:]:
                result.append("| " + " | ".join(row) + " |")

            result.append("")

    return "\n".join(result)


def convert_pptx_to_markdown(source: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        raise RuntimeError(
            "Missing python-pptx. Install it with: pip install python-pptx"
        )

    presentation = Presentation(source)
    result = []

    for slide_number, slide in enumerate(presentation.slides, start=1):
        result.append(f"# Slide {slide_number}")
        result.append("")

        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue

            text = shape.text.strip()

            if text:
                result.append(text)
                result.append("")

    return "\n".join(result)


PDF_TABLE_TIMEOUT_SECONDS = 30


class PdfTableTimeoutError(Exception):
    pass


def extract_pdf_tables(page, timeout=PDF_TABLE_TIMEOUT_SECONDS) -> list:
    if not hasattr(page, "find_tables"):
        return []

    # find_tables() is a heuristic layout-analysis algorithm with no
    # guaranteed runtime bound - on certain real-world pages (dense
    # text, unusual layouts) it can run long enough to hang an
    # unattended batch of thousands of files. Run it on a worker thread
    # so we can give up after `timeout` instead of blocking forever.
    executor = concurrent.futures.ThreadPoolExecutor(max_workers=1)

    try:
        future = executor.submit(lambda: page.find_tables().tables)
        found_tables = future.result(timeout=timeout)
    except concurrent.futures.TimeoutError:
        raise PdfTableTimeoutError(
            f"table detection timed out after {timeout}s"
        )
    finally:
        # Don't wait for an abandoned call to finish - it may never
        # return. The thread is left to finish on its own in the
        # background; the caller must not touch this page's document
        # again afterward (see convert_pdf_to_markdown).
        executor.shutdown(wait=False)

    def cell_to_string(value):
        if value is None:
            return ""
        return str(value).replace("\n", " ").replace("|", r"\|")

    tables = []

    for table in found_tables:
        rows = table.extract()

        if not rows:
            continue

        header = [cell_to_string(cell) for cell in rows[0]]
        lines = [
            "| " + " | ".join(header) + " |",
            "| " + " | ".join("---" for _ in header) + " |",
        ]

        for row in rows[1:]:
            lines.append(
                "| "
                + " | ".join(cell_to_string(cell) for cell in row)
                + " |"
            )

        tables.append(lines)

    return tables


TESSERACT_COMMON_PATHS = (
    r"C:\Program Files\Tesseract-OCR\tesseract.exe",
    r"C:\Program Files (x86)\Tesseract-OCR\tesseract.exe",
)


def locate_tesseract():
    """
    Returns a path to the Tesseract executable if it's installed but not
    on PATH, so pytesseract can be pointed at it directly. A fresh
    winget/installer PATH entry doesn't reach an already-open terminal,
    so relying on PATH alone means "install Tesseract" doesn't actually
    fix OCR until the user finds and restarts their shell. Returns None
    if Tesseract is already on PATH (nothing to do), or isn't found in
    any known common location either.
    """
    if shutil.which("tesseract"):
        return None

    for candidate in TESSERACT_COMMON_PATHS:
        if Path(candidate).is_file():
            return candidate

    return None


def ocr_page(page) -> str:
    try:
        import pytesseract
        from PIL import Image
    except ImportError:
        raise RuntimeError(
            "Missing OCR dependencies. Install them with: "
            "pip install pytesseract Pillow"
        )

    tesseract_path = locate_tesseract()
    if tesseract_path:
        pytesseract.pytesseract.tesseract_cmd = tesseract_path

    pixmap = page.get_pixmap(dpi=300)
    image = Image.frombytes(
        "RGB", (pixmap.width, pixmap.height), pixmap.samples
    )

    try:
        return pytesseract.image_to_string(image).strip()
    except pytesseract.TesseractNotFoundError as exc:
        raise RuntimeError(
            "Tesseract OCR engine not found. Install it and make sure "
            "it is on PATH. On Windows: "
            "winget install UB-Mannheim.TesseractOCR"
        ) from exc


def convert_pdf_to_markdown(source: Path) -> str:
    try:
        import pymupdf as fitz
    except ImportError:
        raise RuntimeError(
            "Missing pymupdf. Install it with: pip install pymupdf"
        )

    document = fitz.open(source)
    result = []
    timed_out = False

    try:
        for page_number, page in enumerate(document, start=1):
            result.append(f"# Page {page_number}")
            result.append("")

            text = page.get_text().strip()

            if len(text) < 10:
                text = ocr_page(page)

            if text:
                result.append(text)
                result.append("")

            try:
                for table_lines in extract_pdf_tables(page):
                    result.extend(table_lines)
                    result.append("")
            except PdfTableTimeoutError:
                # The worker thread that ran find_tables() is abandoned,
                # not killed, and may still be reading this document in
                # the background - closing it here could race with that
                # thread. Give up on the whole PDF and leave the
                # document to be released once that thread finishes.
                timed_out = True
                raise
    finally:
        if not timed_out:
            document.close()

    return "\n".join(result)


# ------------------------------------------------------------
# File classification
# ------------------------------------------------------------

def is_supported_for_text_processing(path: Path) -> bool:
    return path.suffix.lower() in TEXT_EXTENSIONS


def is_office_file(path: Path) -> bool:
    return path.suffix.lower() in OFFICE_EXTENSIONS


# ------------------------------------------------------------
# Archive handling
# ------------------------------------------------------------

ARCHIVE_SUFFIXES = (
    (".tar.gz", "tar"),
    (".tar.bz2", "tar"),
    (".tar.xz", "tar"),
    (".tgz", "tar"),
    (".tbz2", "tar"),
    (".txz", "tar"),
    (".tar", "tar"),
    (".zip", "zip"),
    (".7z", "7z"),
    (".rar", "rar"),
)


def classify_archive(path: Path) -> Optional[str]:
    name = path.name.lower()

    for suffix, kind in ARCHIVE_SUFFIXES:
        if name.endswith(suffix):
            return kind

    return None


def archive_stem(path: Path) -> str:
    name = path.name

    for suffix, _kind in ARCHIVE_SUFFIXES:
        if name.lower().endswith(suffix):
            stem = name[: -len(suffix)]

            # A member literally named ".zip" (or "..zip", "...zip")
            # strips down to an empty or dot-only stem, which is not a
            # usable path component: Path.with_name("") and
            # with_name(".") raise ValueError and abort the whole run,
            # and with_name("..") silently builds a path that climbs out
            # of the output directory. Fall back to the on-disk name,
            # which is a valid component by construction. (Not
            # path.stem: before Python 3.14 that returns "." for
            # "..zip" and ".." for "...zip" - the very values being
            # guarded against here.)
            if stem and stem not in (".", ".."):
                return stem

            return name

    return path.stem


def is_archive_file(path: Path) -> bool:
    return classify_archive(path) is not None


def extract_archive(source: Path, dest_dir: Path) -> None:
    kind = classify_archive(source)

    if kind == "zip":
        import zipfile

        with zipfile.ZipFile(source) as archive:
            archive.extractall(dest_dir)
        return

    if kind == "tar":
        import tarfile

        with tarfile.open(source) as archive:
            try:
                archive.extractall(dest_dir, filter="data")
            except TypeError:
                # Python < 3.12 does not support the `filter` argument,
                # and extracting unfiltered would let a crafted member
                # ("../../evil") write outside dest_dir - outside the
                # temp tree main() cleans up. Check every member first.
                dest_root = Path(dest_dir).resolve()

                for member in archive.getmembers():
                    member_path = (dest_root / member.name).resolve()

                    if (
                        dest_root != member_path
                        and dest_root not in member_path.parents
                    ):
                        raise ValueError(
                            f"Unsafe path in tar archive: {member.name}"
                        )

                archive.extractall(dest_dir)
        return

    if kind == "7z":
        try:
            import py7zr
        except ImportError:
            raise RuntimeError(
                "Missing py7zr. Install it with: pip install py7zr"
            )

        with py7zr.SevenZipFile(source, mode="r") as archive:
            archive.extractall(dest_dir)
        return

    if kind == "rar":
        try:
            import rarfile
        except ImportError:
            raise RuntimeError(
                "Missing rarfile. Install it with: pip install rarfile"
            )

        try:
            with rarfile.RarFile(source) as archive:
                archive.extractall(dest_dir)
        except rarfile.NeedFirstVolume:
            raise
        except rarfile.Error as exc:
            raise RuntimeError(
                "Cannot extract RAR archive. This usually means the "
                "'unrar' tool is not installed and on PATH, or the "
                f"archive is corrupt/password-protected: {exc}"
            ) from exc
        return

    raise ValueError(f"Not a supported archive: {source}")


# ------------------------------------------------------------
# Progress bar
# ------------------------------------------------------------

def format_duration(seconds: float) -> str:
    if seconds < 1:
        return "<1s"

    seconds = int(seconds)

    hours, remainder = divmod(seconds, 3600)
    minutes, seconds = divmod(remainder, 60)

    if hours:
        return f"{hours}h {minutes:02d}m {seconds:02d}s"

    if minutes:
        return f"{minutes}m {seconds:02d}s"

    return f"{seconds}s"


class ProgressBar:
    def __init__(self, total: int, width: int = 40):
        self.total = total
        self.width = width
        self.current = 0
        self.start_time = time.monotonic()

    def update(self, current: int, current_file: Path):
        self.current = current

        elapsed = time.monotonic() - self.start_time

        if self.current > 0:
            rate = self.current / elapsed
            remaining = (self.total - self.current) / rate
        else:
            remaining = 0

        percent = (
            100.0
            if self.total == 0
            else self.current / self.total * 100
        )

        filled = (
            self.width
            if self.total == 0
            else int(self.width * self.current / self.total)
        )

        bar = "█" * filled + "░" * (self.width - filled)

        filename = str(current_file)

        # Keep the terminal line reasonably compact.
        terminal_width = shutil.get_terminal_size(
            fallback=(120, 20)
        ).columns

        prefix = (
            f"\r[{bar}] "
            f"{percent:6.2f}% "
            f"({self.current}/{self.total}) "
            f"ETA {format_duration(remaining)} "
        )

        available = max(20, terminal_width - len(prefix) - 1)
        filename = filename[-available:]

        print(
            prefix + filename,
            end="",
            flush=True,
        )

    def finish(self):
        print()


# ------------------------------------------------------------
# Scan
# ------------------------------------------------------------

MAX_ARCHIVE_DEPTH = 10


class ScanEntry(NamedTuple):
    source: Path
    relative_destination: Path


def scan_files(
    source_dir: Path,
    output_root: Path,
    temp_dirs: list,
    extraction_errors: list,
    relative_root: Path = None,
    depth: int = 0,
):
    """
    Recursively find all files under source_dir, excluding the output
    directory. Archives are extracted into fresh temp directories (added
    to temp_dirs for the caller to clean up) and their contents are
    folded into the result with a relative_destination that mirrors the
    archive's location, e.g. bundle.zip -> bundle/<path inside zip>.

    An archive that fails to extract, or that would exceed
    MAX_ARCHIVE_DEPTH, is excluded from the result entirely instead of
    raising or being copied in unchanged: its contents could not be
    inspected, so copying it verbatim would put un-anonymized content
    (dictionary terms, PII) into the output. Extraction failures (but
    not depth-limit fallbacks, which are not errors) are appended to
    extraction_errors so the caller can fold them into its own error
    count.
    """

    if relative_root is None:
        relative_root = Path(".")

    entries = []

    for path in sorted(source_dir.rglob("*")):
        if not path.is_file():
            continue

        try:
            path.relative_to(output_root)
            continue
        except ValueError:
            pass

        relative_path = relative_root / path.relative_to(source_dir)

        if is_archive_file(path):
            if depth >= MAX_ARCHIVE_DEPTH:
                # Deliberate policy fallback, not a failure: logged, but
                # never added to extraction_errors. The archive is
                # excluded rather than copied in unchanged, since its
                # contents were never inspected/anonymized.
                print(
                    f"WARNING: archive nesting exceeds depth limit "
                    f"({MAX_ARCHIVE_DEPTH}), excluding from output: {path}",
                    file=sys.stderr,
                )
                continue

            extract_dir = Path(tempfile.mkdtemp(prefix="anonymize_"))

            try:
                extract_archive(path, extract_dir)
            except Exception as exc:
                shutil.rmtree(extract_dir, ignore_errors=True)
                print(
                    f"ERROR: cannot extract {path}: {exc}",
                    file=sys.stderr,
                )
                extraction_errors.append(path)
                continue

            temp_dirs.append(extract_dir)
            nested_root = relative_path.with_name(archive_stem(path))

            entries.extend(
                scan_files(
                    extract_dir,
                    output_root,
                    temp_dirs,
                    extraction_errors,
                    relative_root=nested_root,
                    depth=depth + 1,
                )
            )
            continue

        entries.append(ScanEntry(path, relative_path))

    return entries


def anonymize_relative_path(relative_path: Path, replacements: dict) -> Path:
    """
    Anonymizes every path segment (folder names, and the file name's
    stem) using the same dictionary substitution as file contents -- the
    output directory must not leak dictionary terms through its folder
    or file names even when a file's own content is otherwise
    anonymized (or isn't touched at all, e.g. a binary file copied
    unchanged still lives under an anonymized folder).

    The final segment's extension is preserved untouched: it keeps the
    file recognizable/openable, and it's what process_file's own
    extension-swapping (e.g. .docx -> .md) operates on afterward.
    """
    parts = relative_path.parts

    if not parts:
        return relative_path

    anonymized_parts = [
        anonymize_text(part, replacements) for part in parts[:-1]
    ]

    file_name = parts[-1]
    stem = Path(file_name).stem
    suffix = Path(file_name).suffix
    anonymized_parts.append(anonymize_text(stem, replacements) + suffix)

    return Path(*anonymized_parts)


def anonymize_destinations(entries: list, replacements: dict) -> list:
    """
    Rewrites every entry's relative_destination via
    anonymize_relative_path. Must run before deduplicate_destinations:
    case-insensitive matching means two differently-cased source names
    (BDR/, bdr/) can anonymize to the same destination, and dedup is
    what resolves that collision.
    """
    return [
        entry._replace(
            relative_destination=anonymize_relative_path(
                entry.relative_destination, replacements
            )
        )
        for entry in entries
    ]


def deduplicate_destinations(entries: list) -> list:
    """
    Make every relative_destination in `entries` unique.

    Archive contents are folded into the output under the archive's stem
    (bundle.zip -> bundle/), so an archive can collide with a plain
    sibling directory of the same name, and two archives can share a stem
    (data.zip and data.tar.gz). Without this step the second writer
    silently overwrites the first and files disappear from the output.

    Colliding entries keep their position but get "__2", "__3", ...
    appended to the final path component, and the collision is logged.
    """

    seen = set()
    result = []

    for entry in entries:
        destination = entry.relative_destination

        if destination not in seen:
            seen.add(destination)
            result.append(entry)
            continue

        counter = 2

        while True:
            candidate = destination.with_name(
                f"{destination.stem}__{counter}{destination.suffix}"
            )

            if candidate not in seen:
                break

            counter += 1

        print(
            f"WARNING: output path collision, writing "
            f"{entry.source} to {candidate} instead of {destination}",
            file=sys.stderr,
        )

        seen.add(candidate)
        result.append(entry._replace(relative_destination=candidate))

    return result


# ------------------------------------------------------------
# Process one file
# ------------------------------------------------------------

def process_file(
    source: Path,
    destination: Path,
    replacements: dict,
) -> tuple[str, int]:
    suffix = source.suffix.lower()

    # Text/code files
    if suffix in TEXT_EXTENSIONS:
        try:
            text = source.read_text(encoding="utf-8")
        except UnicodeDecodeError:
            # Keep binary/non-UTF8 files untouched.
            destination.parent.mkdir(parents=True, exist_ok=True)
            shutil.copy2(source, destination)
            return "copied-non-utf8", 0

        text = anonymize_text(text, replacements)
        text, pii_count = redact_pii(text)

        destination.parent.mkdir(parents=True, exist_ok=True)
        destination.write_text(text, encoding="utf-8")

        return "anonymized", pii_count

    # DOCX -> MD
    if suffix == ".docx":
        markdown = convert_docx_to_markdown(source)
        markdown = anonymize_text(markdown, replacements)
        markdown, pii_count = redact_pii(markdown)

        destination = destination.with_suffix(".md")
        destination.parent.mkdir(parents=True, exist_ok=True)
        destination.write_text(markdown, encoding="utf-8")

        return "office", pii_count

    # XLSX -> MD
    if suffix == ".xlsx":
        markdown = convert_xlsx_to_markdown(source)
        markdown = anonymize_text(markdown, replacements)
        markdown, pii_count = redact_pii(markdown)

        destination = destination.with_suffix(".md")
        destination.parent.mkdir(parents=True, exist_ok=True)
        destination.write_text(markdown, encoding="utf-8")

        return "office", pii_count

    # PPTX -> MD
    if suffix == ".pptx":
        markdown = convert_pptx_to_markdown(source)
        markdown = anonymize_text(markdown, replacements)
        markdown, pii_count = redact_pii(markdown)

        destination = destination.with_suffix(".md")
        destination.parent.mkdir(parents=True, exist_ok=True)
        destination.write_text(markdown, encoding="utf-8")

        return "office", pii_count

    # PDF -> MD
    if suffix == ".pdf":
        markdown = convert_pdf_to_markdown(source)
        markdown = anonymize_text(markdown, replacements)
        markdown, pii_count = redact_pii(markdown)

        destination = destination.with_suffix(".md")
        destination.parent.mkdir(parents=True, exist_ok=True)
        destination.write_text(markdown, encoding="utf-8")

        return "pdf", pii_count

    # Unknown/binary file:
    # Copy it unchanged.
    destination.parent.mkdir(parents=True, exist_ok=True)
    shutil.copy2(source, destination)

    return "copied", 0


# ------------------------------------------------------------
# Main
# ------------------------------------------------------------

def main():
    parser = argparse.ArgumentParser(
        description=(
            "Recursively anonymize text/code files and convert "
            "DOCX/XLSX/PPTX files to Markdown."
        )
    )

    parser.add_argument(
        "source",
        help="Source directory",
    )

    parser.add_argument(
        "replacements",
        help="JSON file containing replacement mappings",
    )

    parser.add_argument(
        "--output",
        help=(
            "Optional output directory. "
            "Default: <source-parent>/anonimized/<source-name>"
        ),
    )

    args = parser.parse_args()

    source_dir = Path(args.source).expanduser().resolve()
    replacements_file = Path(args.replacements).expanduser().resolve()

    if not source_dir.exists():
        print(f"ERROR: source directory does not exist: {source_dir}")
        sys.exit(1)

    if not source_dir.is_dir():
        print(f"ERROR: source path is not a directory: {source_dir}")
        sys.exit(1)

    if not replacements_file.exists():
        print(f"ERROR: replacement file does not exist: {replacements_file}")
        sys.exit(1)

    try:
        replacements = load_replacements(replacements_file)
    except Exception as exc:
        print(f"ERROR: cannot load replacement map: {exc}")
        sys.exit(1)

    if args.output:
        output_root = Path(args.output).expanduser().resolve()
    else:
        output_root = (
            source_dir.parent
            / "anonimized"
            / source_dir.name
        )

    # --------------------------------------------------------
    # PASS 1 - scan
    # --------------------------------------------------------

    print()
    print("==============================================")
    print(" ANONYMIZATION")
    print("==============================================")
    print()
    print(f"Source       : {source_dir}")
    print(f"Output       : {output_root}")
    print(f"Replacements : {len(replacements)}")
    print()
    print("Scanning files...")

    scan_start = time.monotonic()

    temp_dirs = []
    extraction_errors = []

    try:
        files = scan_files(
            source_dir,
            output_root,
            temp_dirs,
            extraction_errors,
        )

        # Anonymize folder/file names before deduplicating: two
        # differently-cased source names (BDR/, bdr/) can now land on
        # the same destination, and dedup is what resolves that.
        files = anonymize_destinations(files, replacements)

        # An extracted archive can land on the same output path as a
        # plain sibling (data.zip next to data/) or as another archive
        # with the same stem. Rename instead of silently overwriting.
        files = deduplicate_destinations(files)

        archives_extracted = len(temp_dirs)

        scan_duration = time.monotonic() - scan_start

        print(
            f"Found {len(files):,} files "
            f"in {format_duration(scan_duration)}."
        )
        print()

        if not files:
            print("Nothing to process.")
            return

        # --------------------------------------------------------
        # PASS 2 - processing
        # --------------------------------------------------------

        output_root.mkdir(
            parents=True,
            exist_ok=True,
        )

        progress = ProgressBar(
            total=len(files)
        )

        processed = 0
        anonymized = 0
        office_converted = 0
        pdf_converted = 0
        copied = 0
        errors = 0
        pii_redacted = 0

        processing_start = time.monotonic()

        for entry in files:
            destination = output_root / entry.relative_destination

            try:
                result, pii_count = process_file(
                    entry.source,
                    destination,
                    replacements,
                )

                pii_redacted += pii_count

                if result == "anonymized":
                    anonymized += 1
                elif result == "office":
                    office_converted += 1
                elif result == "pdf":
                    pdf_converted += 1
                else:
                    copied += 1

            except Exception as exc:
                errors += 1

                print()
                print(
                    f"ERROR: {entry.source}: {exc}",
                    file=sys.stderr,
                )

                # Deliberately not copied unchanged: this tool exists to
                # anonymize PII, and a file that failed to process is
                # exactly the file we could not guarantee is safe. The
                # error is already counted and logged; the file is
                # simply excluded from the output rather than risking a
                # raw, un-anonymized copy landing in it.

            processed += 1

            progress.update(
                processed,
                entry.relative_destination,
            )

        progress.finish()

        duration = time.monotonic() - processing_start

        errors += len(extraction_errors)

        # --------------------------------------------------------
        # Summary
        # --------------------------------------------------------

        print()
        print("==============================================")
        print(" DONE")
        print("==============================================")
        print()
        print(f"Files found       : {len(files):,}")
        print(f"Files processed   : {processed:,}")
        print(f"Archives extracted: {archives_extracted:,}")
        print(f"Text/code         : {anonymized:,}")
        print(f"Office -> Markdown: {office_converted:,}")
        print(f"PDF -> Markdown   : {pdf_converted:,}")
        print(f"Copied unchanged  : {copied:,}")
        print(f"PII fragments     : {pii_redacted:,}")
        print(f"Errors            : {errors:,}")
        print(f"Processing time   : {format_duration(duration)}")
        print()
        print(f"Output directory:")
        print(output_root)
        print()
    finally:
        for temp_dir in temp_dirs:
            shutil.rmtree(temp_dir, ignore_errors=True)


if __name__ == "__main__":
    main()
